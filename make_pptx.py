"""
Generate a new, engaging PowerPoint for:
Ex 1K – Applications of Simultaneous Linear Equations (55 min lesson)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x23, 0x5C)   # deep navy
TEAL   = RGBColor(0x00, 0x97, 0x9C)   # accent teal
GOLD   = RGBColor(0xFF, 0xC0, 0x2E)   # warm gold
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
DGRAY  = RGBColor(0x44, 0x44, 0x44)
GREEN  = RGBColor(0x2E, 0x86, 0x48)
ORANGE = RGBColor(0xE8, 0x57, 0x1A)
PURPLE = RGBColor(0x6A, 0x3D, 0x9A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ═══════════════════════════════════════════════════════════════════════════
# Helper utilities
# ═══════════════════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_para(tf, text, font_size=16, bold=False, color=DGRAY,
             align=PP_ALIGN.LEFT, italic=False, space_before=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def add_label_box(slide, label, l, t, w, h, bg=TEAL, fg=WHITE, fs=13, bold=True):
    add_rect(slide, l, t, w, h, fill=bg)
    add_text(slide, label, l+0.05, t+0.03, w-0.1, h-0.06,
             font_size=fs, bold=bold, color=fg, align=PP_ALIGN.CENTER)

def slide_bg(slide, color=LGRAY):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.35, fill=NAVY)
    add_text(slide, title, 0.25, 0.1, 10, 0.75,
             font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.25, 0.78, 10, 0.5,
                 font_size=16, bold=False, color=TEAL, align=PP_ALIGN.LEFT)

def accent_bar(slide, l, t, w, h=0.06, color=GOLD):
    add_rect(slide, l, t, w, h, fill=color)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
# Full navy background
add_rect(slide, 0, 0, 13.33, 7.5, fill=NAVY)
# Teal accent stripe
add_rect(slide, 0, 5.5, 13.33, 0.12, fill=TEAL)
# Gold geometric accent
add_rect(slide, 10.5, 0, 2.83, 7.5, fill=RGBColor(0x1E, 0x2D, 0x72))
add_rect(slide, 10.7, 0.8, 2.2, 0.08, fill=GOLD)
add_rect(slide, 10.7, 6.4, 2.2, 0.08, fill=GOLD)

# Topic tag
add_label_box(slide, "YEAR 10 MATHEMATICS", 0.4, 0.5, 3.2, 0.45,
              bg=TEAL, fg=WHITE, fs=13)
# Main title
add_text(slide, "Applications of", 0.4, 1.2, 9.5, 1.0,
         font_size=48, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "Simultaneous", 0.4, 2.0, 9.5, 1.1,
         font_size=60, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
add_text(slide, "Linear Equations", 0.4, 2.95, 9.5, 1.0,
         font_size=48, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
# Subtitle line
add_text(slide, "Ex 1K  |  55-Minute Lesson", 0.4, 4.1, 9, 0.5,
         font_size=20, bold=False, color=LGRAY, align=PP_ALIGN.LEFT)
# Lesson tags
tags = [("📘 APPLY", 0.4), ("📊 SOLVE", 2.2), ("🌍 CONNECT", 4.0)]
for lbl, lx in tags:
    add_label_box(slide, lbl, lx, 4.85, 1.55, 0.42,
                  bg=RGBColor(0x00, 0x6E, 0x73), fg=WHITE, fs=12)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Lesson Roadmap / 55-min plan
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, LGRAY)
header_bar(slide, "Today's Lesson Roadmap", "55 minutes — where we're headed")

stages = [
    ("⚡ WARM UP",      "0–5 min",   "Quick puzzle to activate prior knowledge", TEAL),
    ("🎯 INTENTIONS",  "5–10 min",  "What we'll learn and why it matters",       PURPLE),
    ("📖 LEARN",       "10–25 min", "Key steps + two worked examples",            NAVY),
    ("🏋 PRACTICE",    "25–45 min", "Graduated exercises with scaffolding",       GREEN),
    ("🎮 CHALLENGE",   "45–52 min", "Real-world problem-solving race",            ORANGE),
    ("🪞 REFLECT",     "52–55 min", "Exit ticket + learning check",               RGBColor(0x8B, 0x00, 0x8B)),
]
for i, (title, time, desc, col) in enumerate(stages):
    lx = 0.35 + i * 2.1
    add_rect(slide, lx, 1.55, 1.85, 3.2, fill=col)
    add_rect(slide, lx, 1.55, 1.85, 0.55, fill=RGBColor(
        max(col[0]-30,0), max(col[1]-30,0), max(col[2]-30,0)))
    add_text(slide, title, lx+0.05, 1.57, 1.75, 0.5,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, time, lx+0.05, 2.13, 1.75, 0.4,
             font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, desc, lx+0.08, 2.55, 1.7, 1.9,
             font_size=11, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
# Bottom note
add_text(slide, "💡 Interactive website open throughout — earn points as you go!",
         0.5, 6.5, 12, 0.55, font_size=15, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Warm-Up (5 min)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, RGBColor(0xFF, 0xF8, 0xE7))
header_bar(slide, "⚡ Warm-Up  |  5 Minutes", "Decode the mystery amounts!")

add_rect(slide, 0.3, 1.5, 8.0, 4.8, fill=WHITE)
accent_bar(slide, 0.3, 1.5, 8.0, 0.07, color=GOLD)

txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(7.6), Inches(4.5))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "🧩  The Snack Bar Problem"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = NAVY

add_para(tf, "At the school canteen:", 15, italic=True, color=DGRAY, space_before=10)
add_para(tf, "   • 2 pies + 3 drinks cost $13.00", 17, color=DGRAY)
add_para(tf, "   • 4 pies + 1 drink cost $15.00", 17, color=DGRAY)
add_para(tf, "", 10)
add_para(tf, "Can you figure out the price of one pie and one drink?", 16,
         bold=True, color=NAVY)
add_para(tf, "", 8)
add_para(tf, "💬 Discuss with your neighbour for 2 minutes, then share!", 14,
         italic=True, color=TEAL)

# Hint box
add_rect(slide, 8.6, 1.5, 4.4, 2.2, fill=TEAL)
add_text(slide, "💡 Hint", 8.7, 1.55, 4.2, 0.45,
         font_size=16, bold=True, color=WHITE)
add_text(slide,
         "Let p = price of a pie\nLet d = price of a drink\n\n"
         "Write TWO equations\nand solve!",
         8.7, 2.0, 4.2, 1.6, font_size=15, color=WHITE)

# Answer reveal box
add_rect(slide, 8.6, 3.9, 4.4, 2.4, fill=NAVY)
add_text(slide, "✅ Answer (reveal after!)", 8.7, 3.95, 4.2, 0.45,
         font_size=13, bold=True, color=GOLD)
add_text(slide,
         "2p + 3d = 13  … (1)\n4p +  d = 15  … (2)\n\n"
         "Solving: p = $3.00,  d = $2.33",
         8.7, 4.42, 4.2, 1.75, font_size=14, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Learning Intentions
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, LGRAY)
header_bar(slide, "🎯 Learning Intentions", "By the end of this lesson you will be able to…")

intentions = [
    ("1", "Set up equations",
     "Translate a real-world worded problem into a pair of linear equations by carefully defining variables."),
    ("2", "Choose your method",
     "Select the most efficient solving method (substitution or elimination) based on the structure of the equations."),
    ("3", "Solve & interpret",
     "Solve simultaneously and communicate the answer back in the context of the original problem."),
    ("4", "Spot special cases",
     "Identify when two equations represent parallel lines (no solution) or the same line (infinite solutions)."),
]
for i, (num, short, detail) in enumerate(intentions):
    lx = 0.35 if i % 2 == 0 else 6.85
    ty = 1.7 if i < 2 else 4.2
    add_rect(slide, lx, ty, 6.1, 2.1, fill=WHITE)
    accent_bar(slide, lx, ty, 6.1, 0.08, color=TEAL if i % 2 == 0 else GOLD)
    # Number circle
    add_rect(slide, lx+0.1, ty+0.25, 0.55, 0.55, fill=NAVY)
    add_text(slide, num, lx+0.1, ty+0.24, 0.55, 0.55,
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, short, lx+0.78, ty+0.15, 5.0, 0.5,
             font_size=18, bold=True, color=NAVY)
    add_text(slide, detail, lx+0.15, ty+0.72, 5.8, 1.2,
             font_size=14, color=DGRAY, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 – The 4-Step Method
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, RGBColor(0xE8, 0xF4, 0xFD))
header_bar(slide, "📐 The 4-Step Method", "A reliable strategy for every worded problem")

steps = [
    (TEAL,   "STEP 1", "Define Variables",
     'Choose letters that make sense.\nWrite them down clearly.\nE.g.  "Let a = number of adults"'),
    (GREEN,  "STEP 2", "Form Equations",
     'Read carefully — each fact gives you ONE equation.\nCheck: you need exactly 2 equations\nfor 2 unknowns.'),
    (ORANGE, "STEP 3", "Solve Simultaneously",
     'Pick elimination or substitution.\nShow ALL working — method marks matter!\nLabel each equation (1) and (2).'),
    (PURPLE, "STEP 4", "Answer in Context",
     'Write a sentence using the original wording.\nInclude units (e.g. $, kg, hours).\nSanity-check: does the answer make sense?'),
]

for i, (col, step, title, body) in enumerate(steps):
    lx = 0.3 + i * 3.2
    # Card
    add_rect(slide, lx, 1.55, 2.9, 5.2, fill=WHITE)
    add_rect(slide, lx, 1.55, 2.9, 0.9, fill=col)
    # Step label
    add_text(slide, step, lx+0.05, 1.57, 2.8, 0.42,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrow connector (not last)
    if i < 3:
        add_rect(slide, lx+2.92, 3.8, 0.25, 0.18, fill=col)
        # Simple arrow text
        add_text(slide, "▶", lx+2.93, 3.75, 0.22, 0.28,
                 font_size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    # Title
    add_text(slide, title, lx+0.1, 2.5, 2.7, 0.52,
             font_size=17, bold=True, color=col, align=PP_ALIGN.CENTER)
    # Body
    add_text(slide, body, lx+0.15, 3.08, 2.65, 3.4,
             font_size=13, color=DGRAY, wrap=True)

# Bottom tip
add_rect(slide, 0.3, 6.9, 12.7, 0.42, fill=NAVY)
add_text(slide, "💡 Pro tip: underline the key numbers and circle the unknowns as you read the question!",
         0.5, 6.92, 12.5, 0.38, font_size=13, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Worked Example 1 (Question)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, LGRAY)
header_bar(slide, "📖 Worked Example 1", "Setting up and solving — ticket sales")

# Question box
add_rect(slide, 0.3, 1.55, 8.5, 2.8, fill=WHITE)
accent_bar(slide, 0.3, 1.55, 8.5, 0.08, color=GOLD)

txb = slide.shapes.add_textbox(Inches(0.45), Inches(1.7), Inches(8.2), Inches(2.5))
txb.word_wrap = True
tf = txb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run()
r.text = "🎟️  The School Concert"
r.font.size = Pt(19); r.font.bold = True; r.font.color.rgb = NAVY

add_para(tf,
    "Adult tickets cost $12 and student tickets cost $7. "
    "A total of 350 tickets were sold, raising $3150.", 15, color=DGRAY, space_before=8)
add_para(tf, "How many adult tickets and how many student tickets were sold?",
         16, bold=True, color=NAVY, space_before=6)

# Step labels on right
step_labels = [
    (TEAL,   "STEP 1\nDefine Variables", 1.62),
    (GREEN,  "STEP 2\nForm Equations",   2.55),
    (ORANGE, "STEP 3\nSolve",            3.85),
    (PURPLE, "STEP 4\nAnswer",           5.5),
]
for col, lbl, ty in step_labels:
    add_rect(slide, 9.05, ty, 1.5, 0.7, fill=col)
    add_text(slide, lbl, 9.1, ty+0.03, 1.4, 0.65,
             font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrow line
    add_rect(slide, 8.85, ty+0.3, 0.22, 0.06, fill=col)

# Scaffold working area
add_rect(slide, 0.3, 4.5, 12.7, 2.7, fill=WHITE)
accent_bar(slide, 0.3, 4.5, 12.7, 0.07, color=NAVY)
add_text(slide, "✍️  Your turn — set it up!  (Try before we work through it together)",
         0.45, 4.56, 12.5, 0.45, font_size=14, bold=True, color=NAVY)
# Lined areas
for i in range(4):
    add_rect(slide, 0.45, 5.12 + i*0.44, 12.4, 0.03, fill=RGBColor(0xCC,0xCC,0xCC))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Worked Example 1 (Solution)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, LGRAY)
header_bar(slide, "📖 Worked Example 1 — Solution", "")

cols = [TEAL, GREEN, ORANGE, PURPLE]
steps_data = [
    ("STEP 1 — Define Variables",
     "Let  a  =  number of adult tickets sold\nLet  s  =  number of student tickets sold"),
    ("STEP 2 — Form Equations",
     "Total tickets:    a  +  s  =  350   … (1)\nTotal revenue:  12a  +  7s  =  3150  … (2)"),
    ("STEP 3 — Solve (Elimination)",
     "Multiply (1) by 7:    7a + 7s = 2450    … (3)\n"
     "Subtract (3) from (2):   5a = 700\n"
     "∴  a = 140\nSubstitute into (1):  140 + s = 350  →  s = 210"),
    ("STEP 4 — Answer in Context",
     "140 adult tickets and 210 student tickets were sold.\n"
     "✅ Check: 140 + 210 = 350 ✓   and   12(140) + 7(210) = 1680 + 1470 = 3150 ✓"),
]

for i, (title, body) in enumerate(steps_data):
    lx = 0.3 if i % 2 == 0 else 6.85
    ty = 1.55 if i < 2 else 4.05
    add_rect(slide, lx, ty, 6.1, 2.15, fill=WHITE)
    add_rect(slide, lx, ty, 6.1, 0.52, fill=cols[i])
    add_text(slide, title, lx+0.1, ty+0.07, 5.9, 0.42,
             font_size=14, bold=True, color=WHITE)
    add_text(slide, body, lx+0.15, ty+0.62, 5.8, 1.45,
             font_size=14, color=DGRAY, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Worked Example 2 (Question)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, RGBColor(0xF0, 0xF8, 0xF0))
header_bar(slide, "📖 Worked Example 2", "A more complex application — break-even")

add_rect(slide, 0.3, 1.55, 12.7, 2.7, fill=WHITE)
accent_bar(slide, 0.3, 1.55, 12.7, 0.08, color=TEAL)

txb = slide.shapes.add_textbox(Inches(0.45), Inches(1.7), Inches(12.3), Inches(2.4))
txb.word_wrap = True
tf = txb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run()
r.text = "🏭  The Cupcake Business"
r.font.size = Pt(19); r.font.bold = True; r.font.color.rgb = NAVY

add_para(tf,
    "Emma starts a cupcake business. She spends $240 on equipment (fixed cost) "
    "and $1.50 to make each cupcake. She sells each cupcake for $4.50.",
    15, color=DGRAY, space_before=8)
add_para(tf,
    "(a)  Write equations for Emma's total Cost (C) and total Revenue (R) "
    "in terms of n, the number of cupcakes.",
    15, bold=True, color=NAVY, space_before=6)
add_para(tf,
    "(b)  Find the break-even point — how many cupcakes must she sell?",
    15, bold=True, color=NAVY, space_before=4)

# Scaffold grid
for col_x, col_label, col_color in [(0.3,"Cost Equation", TEAL),(6.65,"Revenue Equation", GREEN)]:
    add_rect(slide, col_x, 4.35, 6.1, 2.85, fill=WHITE)
    add_rect(slide, col_x, 4.35, 6.1, 0.48, fill=col_color)
    add_text(slide, col_label, col_x+0.1, 4.37, 5.9, 0.44,
             font_size=14, bold=True, color=WHITE)
    for i in range(4):
        add_rect(slide, col_x+0.15, 4.95+i*0.5, 5.8, 0.03,
                 fill=RGBColor(0xCC,0xCC,0xCC))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Worked Example 2 (Solution)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, RGBColor(0xF0, 0xF8, 0xF0))
header_bar(slide, "📖 Worked Example 2 — Solution", "Break-even analysis")

# Part (a)
add_rect(slide, 0.3, 1.55, 5.9, 2.5, fill=WHITE)
add_rect(slide, 0.3, 1.55, 5.9, 0.52, fill=TEAL)
add_text(slide, "(a)  Equations", 0.4, 1.59, 5.7, 0.44,
         font_size=15, bold=True, color=WHITE)
add_text(slide,
    "Cost:     C  =  1.5n  +  240\n\n"
    "Revenue:  R  =  4.5n\n\n"
    "(n = number of cupcakes sold)",
    0.45, 2.18, 5.7, 1.75, font_size=16, color=DGRAY)

# Part (b)
add_rect(slide, 6.55, 1.55, 6.45, 2.5, fill=WHITE)
add_rect(slide, 6.55, 1.55, 6.45, 0.52, fill=GREEN)
add_text(slide, "(b)  Break-even: set C = R", 6.65, 1.59, 6.2, 0.44,
         font_size=15, bold=True, color=WHITE)
add_text(slide,
    "1.5n + 240  =  4.5n\n"
    "240  =  3n\n"
    "n  =  80 cupcakes\n\n"
    "∴ Emma must sell 80 cupcakes to break even.",
    6.65, 2.18, 6.2, 1.75, font_size=16, color=DGRAY)

# Graph description
add_rect(slide, 0.3, 4.2, 12.7, 3.0, fill=NAVY)
add_text(slide, "📊 What does this look like graphically?", 0.5, 4.25, 12.5, 0.5,
         font_size=17, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(slide,
    "• The Cost line starts at (0, 240) — fixed cost — and rises with gradient 1.5\n"
    "• The Revenue line passes through the origin with gradient 4.5\n"
    "• They intersect at the point (80, 360) — the break-even point\n"
    "• For n < 80: Cost > Revenue → LOSS       For n > 80: Revenue > Cost → PROFIT",
    0.5, 4.82, 12.5, 2.2, font_size=15, color=WHITE, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Practice Time
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, LGRAY)
header_bar(slide, "🏋 Practice Time  |  20 Minutes", "Graduated exercises — choose your level!")

levels = [
    (GREEN, "🌱 FOUNDATION\n(Q1–5)",
     "Highly scaffolded problems with equation frames provided. "
     "Focus on setting up equations correctly and practising elimination."),
    (TEAL, "📘 STANDARD\n(Q4–9)",
     "Mixed worded problems across different contexts. "
     "Choose your own method. Full working required."),
    (ORANGE, "🔥 ADVANCED\n(Q8–14)",
     "Multi-step problems, break-even scenarios, and proof questions. "
     "Extend to parallel/perpendicular line analysis."),
]

for i, (col, title, desc) in enumerate(levels):
    lx = 0.3 + i * 4.3
    add_rect(slide, lx, 1.6, 4.0, 4.5, fill=WHITE)
    add_rect(slide, lx, 1.6, 4.0, 1.1, fill=col)
    add_text(slide, title, lx+0.1, 1.65, 3.8, 1.0,
             font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, lx+0.15, 2.82, 3.7, 3.1,
             font_size=14, color=DGRAY, wrap=True)

# Tips row
add_rect(slide, 0.3, 6.2, 12.7, 1.1, fill=NAVY)
add_text(slide,
    "📝 Show ALL working    |    ✅ Write final answers as sentences    |    "
    "🤝 Ask your neighbour before asking me    |    🎮 Log answers on the website to earn XP!",
    0.5, 6.28, 12.5, 0.9, font_size=13, bold=False, color=WHITE,
    align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 – Special Cases: Parallel & Same Line
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, RGBColor(0xF5, 0xEC, 0xFF))
header_bar(slide, "⚠️ Special Cases", "When things don't work out as expected…")

cases = [
    (TEAL,   "✅ ONE Solution",
     "Lines intersect at exactly one point.\n\n"
     "Gradients are DIFFERENT.\n\n"
     "e.g.  y = 2x + 1\n       y = x + 4\n\n"
     "Solve to find the unique (x, y).",
     "NORMAL CASE"),
    (ORANGE, "🚫 NO Solution",
     "Lines are PARALLEL — they never meet.\n\n"
     "Same gradient, different y-intercept.\n\n"
     "e.g.  y = 3x + 2\n       y = 3x − 5\n\n"
     "Elimination gives: 0 = 7  (impossible!)",
     "PARALLEL LINES"),
    (PURPLE, "∞  INFINITE Solutions",
     "Lines are IDENTICAL — they sit on top of each other.\n\n"
     "Same gradient AND same y-intercept.\n\n"
     "e.g.  2y = 4x + 6\n       y = 2x + 3\n\n"
     "Elimination gives: 0 = 0  (always true!)",
     "SAME LINE"),
]

for i, (col, title, body, tag) in enumerate(cases):
    lx = 0.3 + i * 4.3
    add_rect(slide, lx, 1.55, 4.0, 5.3, fill=WHITE)
    add_rect(slide, lx, 1.55, 4.0, 0.8, fill=col)
    add_text(slide, tag, lx+0.1, 1.57, 3.8, 0.35,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, lx+0.1, 1.9, 3.8, 0.44,
             font_size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, body, lx+0.15, 2.45, 3.75, 4.2,
             font_size=14, color=DGRAY, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 – Real-World Connections
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, RGBColor(0xE8, 0xF4, 0xFD))
header_bar(slide, "🌍 Where Is This Used in the Real World?",
           "Simultaneous equations are everywhere!")

contexts = [
    ("💰", "Economics",     "Finding equilibrium where supply meets demand."),
    ("🏗️", "Engineering",   "Balancing forces in structures and circuits."),
    ("🧪", "Chemistry",     "Mixing solutions to hit a target concentration."),
    ("📈", "Business",      "Break-even analysis and profit optimisation."),
    ("🏥", "Medicine",      "Calculating drug dosages across compartments."),
    ("🎮", "Game Design",   "Balancing character stats and resource systems."),
]

for i, (icon, field, desc) in enumerate(contexts):
    col = i % 3
    row = i // 3
    lx = 0.4 + col * 4.25
    ty = 1.65 + row * 2.35
    add_rect(slide, lx, ty, 3.9, 2.05, fill=WHITE)
    add_rect(slide, lx, ty, 3.9, 0.07, fill=TEAL if row==0 else GOLD)
    add_text(slide, icon, lx+0.1, ty+0.12, 0.7, 0.65,
             font_size=30, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, field, lx+0.8, ty+0.18, 2.9, 0.48,
             font_size=17, bold=True, color=NAVY)
    add_text(slide, desc, lx+0.15, ty+0.72, 3.65, 1.1,
             font_size=13, color=DGRAY, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 – Challenge Round (Gamification)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(slide, 0, 0, 13.33, 0.08, fill=GOLD)
add_rect(slide, 0, 7.42, 13.33, 0.08, fill=GOLD)

add_text(slide, "🎮  CHALLENGE ROUND", 0.5, 0.4, 12.3, 0.8,
         font_size=38, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(slide, "7 minutes  —  solve as many as you can  —  earn XP on the website!",
         0.5, 1.1, 12.3, 0.5, font_size=16, color=LGRAY, align=PP_ALIGN.CENTER)

challenges = [
    ("⭐ 10 XP", "QUICK FIRE",
     "The sum of two numbers is 56.\nTheir difference is 14.\nFind both numbers."),
    ("⭐⭐ 20 XP", "REAL WORLD",
     "A phone plan charges $25/month + $0.10/text.\nAnother charges $15/month + $0.25/text.\nFor how many texts are they equal?"),
    ("⭐⭐⭐ 30 XP", "BOSS LEVEL",
     "Two cars leave cities 480 km apart at the same time.\nCar A travels at 90 km/h, Car B at 70 km/h.\nWhen and where do they meet?"),
]

for i, (xp, title, problem) in enumerate(challenges):
    lx = 0.4 + i * 4.25
    add_rect(slide, lx, 1.8, 3.9, 5.2, fill=RGBColor(0x1E, 0x2D, 0x72))
    add_rect(slide, lx, 1.8, 3.9, 0.07, fill=TEAL)
    add_text(slide, xp, lx+0.1, 1.88, 3.7, 0.42,
             font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, title, lx+0.1, 2.32, 3.7, 0.5,
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, lx+0.15, 2.88, 3.6, 0.03, fill=TEAL)
    add_text(slide, problem, lx+0.15, 3.0, 3.65, 3.6,
             font_size=15, color=LGRAY, wrap=True)

add_text(slide, "🏆 Top 3 scorers on the leaderboard win a bonus prize!",
         0.5, 7.1, 12.3, 0.3, font_size=13, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 – Exit Ticket / Reflect
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide, LGRAY)
header_bar(slide, "🪞 Reflect & Exit Ticket", "3 minutes — show what you know")

# Traffic light self-assessment
add_rect(slide, 0.3, 1.55, 4.5, 5.55, fill=WHITE)
accent_bar(slide, 0.3, 1.55, 4.5, 0.07, color=NAVY)
add_text(slide, "Self-Assessment\nTraffic Light", 0.4, 1.65, 4.3, 0.75,
         font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

lights = [
    (GREEN,  "🟢 Green  — Got it!",
     "I can set up AND solve worded simultaneous equation problems independently."),
    (GOLD,   "🟡 Yellow — Nearly!",
     "I can set up the equations but need more practice with the solving methods."),
    (RGBColor(0xD0,0x32,0x2A), "🔴 Red — Need help!",
     "I'm still unsure about how to start worded problems — more practice needed."),
]
for i, (col, label, desc) in enumerate(lights):
    ty = 2.5 + i * 1.55
    add_rect(slide, 0.4, ty, 4.2, 1.35, fill=RGBColor(0xFA,0xFA,0xFA))
    add_rect(slide, 0.4, ty, 0.18, 1.35, fill=col)
    add_text(slide, label, 0.65, ty+0.1, 3.8, 0.4,
             font_size=13, bold=True, color=DGRAY)
    add_text(slide, desc, 0.65, ty+0.52, 3.8, 0.75,
             font_size=12, color=DGRAY, wrap=True)

# Exit ticket question
add_rect(slide, 5.05, 1.55, 8.0, 5.55, fill=WHITE)
accent_bar(slide, 5.05, 1.55, 8.0, 0.07, color=GOLD)
add_text(slide, "📝 Exit Ticket", 5.15, 1.64, 7.8, 0.5,
         font_size=18, bold=True, color=NAVY)
add_text(slide,
    "On your mini whiteboard (or worksheet):\n\n"
    "A jar contains 20-cent and 50-cent coins.\n"
    "There are 30 coins worth $12.00 in total.\n\n"
    "Find the number of each type of coin.",
    5.15, 2.22, 7.8, 2.2, font_size=16, color=DGRAY, wrap=True)
# Lines
for i in range(5):
    add_rect(slide, 5.15, 4.55+i*0.44, 7.7, 0.03,
             fill=RGBColor(0xCC,0xCC,0xCC))
# Homework note
add_rect(slide, 5.05, 6.5, 8.0, 0.55, fill=NAVY)
add_text(slide, "📚 Homework: Exercise 1K — see working programs on the class portal",
         5.15, 6.54, 7.8, 0.45, font_size=12, bold=True, color=GOLD)

prs.save("/home/user/mathsteaching/Ex 1K - Applications (New Engaging Lesson).pptx")
print("PowerPoint saved successfully!")
